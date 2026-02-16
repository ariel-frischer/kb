"""Text extraction from various document formats.

Central registry mapping file extensions to extractor functions.
Stdlib-only formats are always available; optional formats require extra deps.
"""

import re
import zipfile
import xml.etree.ElementTree as ET
from email import policy
from email.parser import BytesParser
from html.parser import HTMLParser
from pathlib import Path


# ---------------------------------------------------------------------------
# Optional dependency probes
# ---------------------------------------------------------------------------

try:
    import pymupdf  # noqa: F401

    _PYMUPDF = True
except ImportError:
    _PYMUPDF = False

try:
    import docx as _docx  # noqa: F401

    _DOCX = True
except ImportError:
    _DOCX = False

try:
    import pptx as _pptx  # noqa: F401

    _PPTX = True
except ImportError:
    _PPTX = False

try:
    import openpyxl as _openpyxl  # noqa: F401

    _OPENPYXL = True
except ImportError:
    _OPENPYXL = False

try:
    from striprtf.striprtf import rtf_to_text as _rtf_to_text  # noqa: F401

    _STRIPRTF = True
except ImportError:
    _STRIPRTF = False


# ---------------------------------------------------------------------------
# Extractors — each returns plain text from a file path
# ---------------------------------------------------------------------------


def _read_text(path: Path) -> str:
    return path.read_text(errors="replace")


class _HTMLStripper(HTMLParser):
    """Minimal HTML→text converter using stdlib."""

    def __init__(self) -> None:
        super().__init__()
        self._parts: list[str] = []
        self._skip = False

    def handle_starttag(self, tag: str, attrs: list) -> None:
        if tag in ("script", "style"):
            self._skip = True

    def handle_endtag(self, tag: str) -> None:
        if tag in ("script", "style"):
            self._skip = False
        if tag in ("p", "div", "br", "li", "tr", "h1", "h2", "h3", "h4", "h5", "h6"):
            self._parts.append("\n")

    def handle_data(self, data: str) -> None:
        if not self._skip:
            self._parts.append(data)

    def get_text(self) -> str:
        return re.sub(r"\n{3,}", "\n\n", "".join(self._parts)).strip()


def _strip_html(html: str) -> str:
    s = _HTMLStripper()
    s.feed(html)
    return s.get_text()


def _extract_html(path: Path) -> str:
    return _strip_html(path.read_text(errors="replace"))


def _extract_srt(path: Path) -> str:
    """Extract text from SRT/VTT subtitle files, stripping sequence numbers and timestamps."""
    lines = []
    for line in path.read_text(errors="replace").splitlines():
        line = line.strip()
        # skip blank, sequence numbers, timestamps, VTT header
        if not line:
            continue
        if re.match(r"^\d+$", line):
            continue
        if re.match(r"[\d:.,\->]+\s*-->", line):
            continue
        if line.startswith("WEBVTT"):
            continue
        if re.match(r"^NOTE\b", line):
            continue
        lines.append(line)
    return "\n".join(lines)


def _extract_eml(path: Path) -> str:
    """Extract text from .eml email files."""
    msg = BytesParser(policy=policy.default).parse(path.open("rb"))
    parts = []
    subj = msg.get("subject", "")
    if subj:
        parts.append(f"Subject: {subj}")
    sender = msg.get("from", "")
    if sender:
        parts.append(f"From: {sender}")
    date = msg.get("date", "")
    if date:
        parts.append(f"Date: {date}")
    if parts:
        parts.append("")

    body = msg.get_body(preferencelist=("plain", "html"))
    if body:
        content = body.get_content()
        if body.get_content_type() == "text/html":
            content = _strip_html(content)
        parts.append(content)
    return "\n".join(parts)


def _extract_odt(path: Path) -> str:
    """Extract text from OpenDocument Text files."""
    ns = {"text": "urn:oasis:names:tc:opendocument:xmlns:text:1.0"}
    with zipfile.ZipFile(path) as zf:
        tree = ET.fromstring(zf.read("content.xml"))
    paragraphs = []
    for p in tree.iter(f"{{{ns['text']}}}p"):
        text = "".join(p.itertext())
        if text.strip():
            paragraphs.append(text)
    return "\n\n".join(paragraphs)


def _extract_ods(path: Path) -> str:
    """Extract text from OpenDocument Spreadsheet files."""
    ns = {
        "table": "urn:oasis:names:tc:opendocument:xmlns:table:1.0",
        "text": "urn:oasis:names:tc:opendocument:xmlns:text:1.0",
    }
    with zipfile.ZipFile(path) as zf:
        tree = ET.fromstring(zf.read("content.xml"))
    rows = []
    for row in tree.iter(f"{{{ns['table']}}}table-row"):
        cells = []
        for cell in row.iter(f"{{{ns['table']}}}table-cell"):
            text = "".join(cell.itertext()).strip()
            cells.append(text)
        line = "\t".join(cells).strip()
        if line:
            rows.append(line)
    return "\n".join(rows)


def _extract_odp(path: Path) -> str:
    """Extract text from OpenDocument Presentation files."""
    ns = {"text": "urn:oasis:names:tc:opendocument:xmlns:text:1.0"}
    with zipfile.ZipFile(path) as zf:
        tree = ET.fromstring(zf.read("content.xml"))
    paragraphs = []
    for p in tree.iter(f"{{{ns['text']}}}p"):
        text = "".join(p.itertext())
        if text.strip():
            paragraphs.append(text)
    return "\n\n".join(paragraphs)


def _extract_epub(path: Path) -> str:
    """Extract text from EPUB files."""
    parts = []
    with zipfile.ZipFile(path) as zf:
        for name in zf.namelist():
            if name.endswith((".xhtml", ".html", ".htm")):
                html = zf.read(name).decode("utf-8", errors="replace")
                text = _strip_html(html)
                if text.strip():
                    parts.append(text)
    return "\n\n".join(parts)


def _extract_pdf(path: Path) -> str:
    """Extract text from PDF using pymupdf."""
    import pymupdf

    doc = pymupdf.open(str(path))
    pages = []
    for page in doc:
        text = page.get_text()
        if text.strip():
            pages.append(text)
    doc.close()
    return "\n\n".join(pages)


def _extract_docx(path: Path) -> str:
    """Extract text from .docx using python-docx."""
    import docx

    doc = docx.Document(str(path))
    return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _extract_pptx(path: Path) -> str:
    """Extract text from .pptx using python-pptx."""
    from pptx import Presentation

    prs = Presentation(str(path))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)
    return "\n\n".join(parts)


def _extract_xlsx(path: Path) -> str:
    """Extract text from .xlsx using openpyxl."""
    import openpyxl

    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    rows = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            line = "\t".join(cells).strip()
            if line:
                rows.append(line)
    wb.close()
    return "\n".join(rows)


def _extract_rtf(path: Path) -> str:
    """Extract text from .rtf using striprtf."""
    from striprtf.striprtf import rtf_to_text

    return rtf_to_text(path.read_text(errors="replace"))


# ---------------------------------------------------------------------------
# Format registry: ext -> (extractor, doc_type, available, install_hint, is_code)
# ---------------------------------------------------------------------------

_FORMATS: dict[str, tuple[..., str, bool, str | None, bool]] = {}


def _register(
    exts: list[str],
    extractor,
    doc_type: str,
    available: bool = True,
    install_hint: str | None = None,
    *,
    is_code: bool = False,
) -> None:
    for ext in exts:
        _FORMATS[ext] = (extractor, doc_type, available, install_hint, is_code)


# Markdown — no deps
_register([".md", ".markdown"], _read_text, "markdown")

# Plain text — no deps
_register(
    [
        ".txt",
        ".text",
        ".rst",
        ".org",
        ".log",
        ".csv",
        ".tsv",
        ".json",
        ".yaml",
        ".yml",
        ".toml",
        ".xml",
        ".ini",
        ".cfg",
        ".tex",
        ".latex",
        ".bib",
        ".nfo",
        ".adoc",
        ".asciidoc",
        ".properties",
    ],
    _read_text,
    "text",
)

# HTML — stdlib
_register([".html", ".htm", ".xhtml"], _extract_html, "html")

# Subtitles — stdlib
_register([".srt", ".vtt"], _extract_srt, "text")

# Email — stdlib
_register([".eml"], _extract_eml, "email")

# OpenDocument — stdlib (zipfile + xml)
_register([".odt"], _extract_odt, "odt")
_register([".ods"], _extract_ods, "ods")
_register([".odp"], _extract_odp, "odp")

# EPUB — stdlib (zipfile + html strip)
_register([".epub"], _extract_epub, "epub")

# PDF — optional
_register([".pdf"], _extract_pdf, "pdf", _PYMUPDF, "pymupdf")

# Office — optional
_register([".docx"], _extract_docx, "docx", _DOCX, "python-docx")
_register([".pptx"], _extract_pptx, "pptx", _PPTX, "python-pptx")
_register([".xlsx"], _extract_xlsx, "xlsx", _OPENPYXL, "openpyxl")

# RTF — optional
_register([".rtf"], _extract_rtf, "rtf", _STRIPRTF, "striprtf")

# Code files — opt-in via index_code config (read as plain text, doc_type="code")
_register(
    [
        # Python
        ".py",
        ".pyi",
        ".pyw",
        # JavaScript / TypeScript
        ".js",
        ".jsx",
        ".ts",
        ".tsx",
        ".mjs",
        ".cjs",
        # Web
        ".css",
        ".scss",
        ".sass",
        ".less",
        ".vue",
        ".svelte",
        # Systems
        ".c",
        ".h",
        ".cpp",
        ".cxx",
        ".cc",
        ".hpp",
        ".hxx",
        # JVM
        ".java",
        ".kt",
        ".kts",
        ".scala",
        ".groovy",
        ".gradle",
        # .NET
        ".cs",
        ".fs",
        ".vb",
        # Go / Rust / Zig
        ".go",
        ".rs",
        ".zig",
        # Ruby / PHP / Perl
        ".rb",
        ".php",
        ".pl",
        ".pm",
        # Shell
        ".sh",
        ".bash",
        ".zsh",
        ".fish",
        ".ps1",
        ".bat",
        ".cmd",
        # Functional / ML
        ".hs",
        ".lhs",
        ".ml",
        ".mli",
        ".ex",
        ".exs",
        ".erl",
        ".clj",
        ".cljs",
        # Lisp
        ".el",
        ".lisp",
        ".cl",
        ".scm",
        ".rkt",
        # Data / Config (code-adjacent)
        ".sql",
        ".graphql",
        ".gql",
        ".proto",
        # Other
        ".r",
        ".R",
        ".jl",
        ".lua",
        ".nim",
        ".dart",
        ".swift",
        ".m",
        ".mm",  # Obj-C
        ".v",
        ".sv",  # Verilog
        ".vhd",
        ".vhdl",  # VHDL
        ".tf",
        ".hcl",  # Terraform
        ".cmake",
        ".mk",
        ".asm",
        ".s",
        # Notebooks / markup that's code-like
        ".ipynb",
    ],
    _read_text,
    "code",
    is_code=True,
)


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------


def extract_text(path: Path, *, include_code: bool = False) -> tuple[str, str] | None:
    """Extract text from a file. Returns (text, doc_type) or None if unsupported/unavailable."""
    ext = path.suffix.lower()
    entry = _FORMATS.get(ext)
    if entry is None:
        return None
    extractor, doc_type, available, _, is_code = entry
    if not available:
        return None
    if is_code and not include_code:
        return None
    return extractor(path), doc_type


def supported_extensions(*, include_code: bool = False) -> set[str]:
    """Extensions that are currently available (deps installed)."""
    return {
        ext
        for ext, (_, _, avail, _, is_code) in _FORMATS.items()
        if avail and (not is_code or include_code)
    }


def unavailable_formats() -> list[tuple[str, str]]:
    """[(ext, install_pkg)] for formats whose optional deps are missing."""
    seen: dict[str, str] = {}
    for ext, (_, _, avail, hint, _) in _FORMATS.items():
        if not avail and hint and hint not in seen.values():
            seen[ext] = hint
    return list(seen.items())
